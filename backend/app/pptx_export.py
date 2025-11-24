# backend/app/pptx_export.py
from io import BytesIO
from pptx import Presentation
from .models import Project

def build_pptx(project: Project) -> BytesIO:
    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = project.title
    subtitle = slide.placeholders[1]
    subtitle.text = project.topic

    # Content slides
    content_layout = prs.slide_layouts[1]
    for section in sorted(project.sections, key=lambda s: s.order):
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = section.title
        body = slide.placeholders[1]
        body.text = section.content or ""

    file_stream = BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    return file_stream
